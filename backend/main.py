# main.py
# Import necessary libraries
import os
import fitz  # PyMuPDF
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
import google.generativeai as genai
from pptx import Presentation
import json
import io
from dotenv import load_dotenv
from pydantic import BaseModel, Field
from typing import List, Literal
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR

from theme_config import THEME_CONFIG


# --- Models ---
class Slide(BaseModel):
    title: str = Field(..., description="The title of the slide (maximum 10 words)")
    bullets: List[str] = Field(None, description="List of bullet points for the slide")
    text_block: str = Field(None, description="A block of text for the slide instead of bullet points")
    detail_level: int = Field(2, description="Level of detail (0-4): 0=very concise, 1=concise, 2=normal, 3=detailed, 4=very detailed")


class SlideContent(BaseModel):
    slides: List[Slide] = Field(..., description="List of slides for the presentation")


class PresentationContent(BaseModel):
    slides: List[Slide] = Field(..., description="List of slides for the presentation")
    theme_type: Literal["color", "background"] = Field(..., description="The type of theme to apply.")
    theme_name: str = Field(..., description="The name of the theme to apply (e.g., 'classic_dark', 'blue_gradient').")


# --- Configuration ---
load_dotenv()
app = FastAPI(title="PDF to Presentation API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

try:
    api_key = os.getenv("GOOGLE_API_KEY")
    if not api_key:
        raise ValueError("GOOGLE_API_KEY environment variable not set.")
    genai.configure(api_key=api_key)
except ValueError as e:
    print(f"ERROR: {e}")


# --- Helper Functions ---

def extract_text_from_pdf(pdf_content: bytes) -> str:
    """Extracts text content from a PDF file's bytes."""
    full_text = ""
    with fitz.open(stream=pdf_content, filetype="pdf") as doc:
        for page in doc:
            full_text += page.get_text()
    print(f"Successfully extracted {len(full_text)} characters from the PDF.")
    return full_text


def generate_slides_content_with_gemini(text: str, detail_level: int = 2) -> List[Slide]:
    """
    Uses Gemini to generate presentation content from text.

    Parameters:
    - text: The text to generate slides from
    - detail_level: Level of detail (0-4): 0=very concise, 1=concise, 2=normal, 3=detailed, 4=very detailed
    """
    print(f"Generating slide content with Gemini at detail level {detail_level}...")
    model = genai.GenerativeModel('gemini-2.0-flash')

    # Map detail level to descriptive text and bullet point count
    detail_descriptions = {
        0: "very concise",
        1: "concise",
        2: "normal",
        3: "detailed",
        4: "very detailed"
    }

    # Map detail level to approximate number of bullet points
    bullet_counts = {
        0: "1-2",
        1: "2-3",
        2: "3-4",
        3: "4-5",
        4: "5-6"
    }

    detail_description = detail_descriptions.get(detail_level, "normal")
    bullet_count = bullet_counts.get(detail_level, "3-4")

    prompt = f"""
    Based on the following text from a report, please generate a {detail_description} summary presentation.
    The output should be a valid JSON object.

    The JSON object must be a single list `[]` containing multiple slide objects {{}}.

    Each slide object must have the following structure:
    {{
        "title": "Slide Title",  // A string for the slide's title (maximum 10 words)

        // EITHER use bullets OR text_block, not both
        "bullets": ["Point 1", "Point 2", ...],  // A list of {bullet_count} bullet points
        // OR
        "text_block": "A paragraph of text that summarizes the key points..."  // A block of text instead of bullet points
    }}

    About 20% of slides should use text_block instead of bullets for variety.
    For bullet point slides, include approximately {bullet_count} bullet points per slide.
    The level of detail should be {detail_description}.

    Summarize the key points and structure them logically for a presentation.

    Here is the text:
    ---
    {text}
    ---
    """
    try:
        response = model.generate_content(prompt)
        response_text = response.text.strip().replace("```json", "").replace("```", "")
        raw_slide_data = json.loads(response_text)

        if not isinstance(raw_slide_data, list):
            raise ValueError("AI response is not a list.")

        # Validate each slide has a title and either bullets or text_block
        for slide in raw_slide_data:
            if "title" not in slide:
                raise ValueError("AI response is missing required key 'title'.")
            if "bullets" not in slide and "text_block" not in slide:
                raise ValueError("AI response is missing both 'bullets' and 'text_block'. At least one is required.")

        # Convert raw dictionaries to Slide objects
        slide_data = []
        for slide in raw_slide_data:
            slide_obj = Slide(
                title=slide["title"],
                detail_level=detail_level
            )

            # Set either bullets or text_block
            if "bullets" in slide and slide["bullets"]:
                slide_obj.bullets = slide["bullets"]
            elif "text_block" in slide and slide["text_block"]:
                slide_obj.text_block = slide["text_block"]
            else:
                # Fallback to empty bullets if neither is provided
                slide_obj.bullets = []

            slide_data.append(slide_obj)

        print("Successfully generated and parsed slide data from Gemini.")
        return slide_data
    except Exception as e:
        print(f"An error occurred during Gemini content generation: {e}")
        raise HTTPException(status_code=500, detail="Failed to generate or parse content from AI.")


def set_slide_background_image(prs, slide, image_path):
    """
    Sets the background of a slide to an image by adding a full-slide picture
    and sending it to the back.

    WARNING: This method uses internal, non-public APIs of python-pptx (e.g., _spTree).
    This means it is subject to breaking changes in future versions of the library
    without prior notice. Use with caution and be prepared to update your code.

    Args:
        prs (pptx.presentation.Presentation): The presentation object.
        slide (pptx.slide.Slide): The slide object to modify.
        image_path (str): The path to the image file.
    """
    left = top = 0
    width = prs.slide_width
    height = prs.slide_height

    # Add the picture to the slide, covering the entire area
    pic = slide.shapes.add_picture(image_path, left, top, width, height)

    # --- START OF INTERNAL API USAGE ---
    # These lines manipulate the internal XML structure (shape tree)
    # to change the Z-order of the picture. This is not part of the
    # public API and might break in future python-pptx versions.
    try:
        # Accessing shape_id forces the shape to be created in the XML tree
        _ = pic.shape_id
        # Remove the picture element from its current position
        slide.shapes._spTree.remove(pic._element)
        # Insert it at index 2 to send it to the back
        # (index 0 and 1 are usually for XML declaration and other core elements)
        slide.shapes._spTree.insert(2, pic._element)
    except AttributeError:
        # Fallback if internal structure changes dramatically (unlikely for _element)
        print("Warning: Could not manipulate Z-order using internal APIs. "
              "Image might not be at the back.")
    # --- END OF INTERNAL API USAGE ---


def create_presentation(slide_data: List[Slide], theme_type: str, theme_name: str) -> io.BytesIO:
    """Creates a PowerPoint presentation from structured data using a predefined theme."""
    print(f"Creating presentation with theme '{theme_name}' (type: {theme_type})")

    prs = Presentation()
    prs.slide_width = 9144000
    prs.slide_height = 5143500

    # Pre-define layouts for clarity and direct access to masters
    title_and_content_layout = prs.slide_layouts[1]  # For bullet points
    title_only_layout = prs.slide_layouts[5]  # For text blocks

    # Define a consistent font size for body text and text blocks
    DEFAULT_BODY_FONT_SIZE = Pt(20)  # You can adjust this value as needed

    # Find the selected theme from the configuration
    theme = None
    if theme_type == "color":
        theme = next((t for t in THEME_CONFIG["colors"] if t["name"] == theme_name), None)
    elif theme_type == "background":
        theme = next((t for t in THEME_CONFIG["backgrounds"] if t["name"] == theme_name), None)

    if not theme:
        print(f"Warning: Theme '{theme_name}' not found. Using default blank presentation.")

    # *** NEW: Apply theme to slide masters if it's a color theme ***
    if theme and theme_type == "color":
        primary_color_rgb = RGBColor.from_string(theme["primaryColor"])
        secondary_color_rgb = RGBColor.from_string(theme["secondaryColor"])
        text_color_rgb = RGBColor.from_string(theme["textColor"])

        # Apply to title_and_content_layout's master
        master_slide_1 = title_and_content_layout.slide_master
        fill = master_slide_1.background.fill
        fill.solid()
        fill.fore_color.rgb = primary_color_rgb

        # Apply to title_only_layout's master
        master_slide_5 = title_only_layout.slide_master
        fill = master_slide_5.background.fill
        fill.solid()
        fill.fore_color.rgb = primary_color_rgb

        # Optional: Set default text colors on masters (e.g., for title placeholder styles)
        # This can be more complex as it involves iterating through text styles on masters.
        # For simplicity, we'll continue to set explicit text color on individual slides
        # where we apply the text, as that gives more direct control.
        # But for title, it's often a good idea to set it on the master level.
        # Example for title on master (might vary based on placeholder index/name):
        # for sh in master_slide_1.shapes:
        #     if sh.has_text_frame and sh.is_placeholder:
        #         if sh.placeholder_format.type == MSO_PLACEHOLDER.TITLE:
        #             sh.text_frame.paragraphs[0].font.color.rgb = secondary_color_rgb
        #             break # Found title

    for slide_info in slide_data:
        slide = None
        title_shape = None

        # Determine which layout to use based on content type
        if slide_info.text_block:
            slide = prs.slides.add_slide(title_only_layout)
        elif slide_info.bullets:
            slide = prs.slides.add_slide(title_and_content_layout)
        else:
            print(
                f"Warning: Slide '{slide_info.title}' has no content (neither text_block nor bullets). Using default Title and Content layout.")
            slide = prs.slides.add_slide(title_and_content_layout)

        # Get title shape regardless of layout (most layouts have a title)
        try:
            title_shape = slide.shapes.title
        except AttributeError:
            print(f"Warning: Slide '{slide_info.title}' layout does not have a title placeholder.")

        # Apply theme specifics (only image background per-slide now)
        if theme:
            if theme_type == "color":
                # Background color is now handled by the master, so nothing to do here.
                # Title text color is still set here for direct control, but could be on master.
                if title_shape:
                    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme["secondaryColor"])

            elif theme_type == "background":
                # Background image must still be set per-slide
                image_path = os.path.join('backgrounds', theme["image_filename"])
                if os.path.exists(image_path):
                    set_slide_background_image(prs, slide, image_path)
                    # Set default text colors for readability on image backgrounds
                    if title_shape:
                        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme["titleColor"])
                else:
                    print(f"Warning: Background image not found at {image_path}. Using default.")

        # Add content to slide
        if title_shape:
            title_shape.text = slide_info.title

        # Determine text color based on theme
        # This will still be determined per-slide as different theme types might imply different text colors
        body_text_color = None
        if theme and (theme_type == 'color' or theme_type == 'background'):
            body_text_color = RGBColor.from_string(theme["textColor"])

        if slide_info.text_block:
            # For "Title Only" layout, we manually add a text box for the content
            left = Inches(1)
            top = Inches(1.75)
            width = Inches(8)
            height = Inches(4)  # Initial height, will auto-adjust

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.clear()

            p = tf.paragraphs[0]
            p.text = slide_info.text_block

            p.font.size = DEFAULT_BODY_FONT_SIZE
            if body_text_color:
                p.font.color.rgb = body_text_color
            p.level = 0

            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            tf.vertical_anchor = MSO_ANCHOR.TOP

        elif slide_info.bullets:
            # For "Title and Content" layout, use the default body placeholder
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()  # Clear existing text (like "Click to add text")

            for i, bullet_text in enumerate(slide_info.bullets):
                if i == 0:  # Use the first paragraph that's already there
                    p = tf.paragraphs[0]
                    p.text = bullet_text
                else:  # Add new paragraphs for subsequent bullets
                    p = tf.add_paragraph()
                    p.text = bullet_text

                p.font.size = DEFAULT_BODY_FONT_SIZE
                if body_text_color:
                    p.font.color.rgb = body_text_color

                p.level = 0

            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    return pptx_stream
# --- API Endpoints ---
@app.post("/api/generate-slide-content")
async def generate_slide_content(
    file: UploadFile = File(...),
    detail_level: int = 2
) -> SlideContent:
    """
    This endpoint receives a PDF file, processes it, and returns the proposed slide content as JSON.

    Parameters:
    - file: The PDF file to process
    - detail_level: Level of detail (0-4): 0=very concise, 1=concise, 2=normal, 3=detailed, 4=very detailed
    """
    print(f"Received file: {file.filename} with detail level: {detail_level}")

    if detail_level < 0 or detail_level > 4:
        raise HTTPException(status_code=400, detail="Detail level must be between 0 and 4")

    if file.content_type != "application/pdf":
        raise HTTPException(status_code=400, detail="Invalid file type. Please upload a PDF.")

    try:
        pdf_contents = await file.read()

        text_content = extract_text_from_pdf(pdf_contents)
        slides = generate_slides_content_with_gemini(text_content, detail_level)

        # Create a Presentation_Content object
        slide_content = SlideContent(slides=slides)

        print("Sending slide content to client.")
        return slide_content
    except Exception as e:
        if isinstance(e, HTTPException):
            raise e
        print(f"An unexpected server error occurred: {e}")
        raise HTTPException(status_code=500, detail=f"An unexpected server error occurred: {str(e)}")


@app.post("/api/generate-presentation")
async def generate_presentation_endpoint(presentation_content: PresentationContent):
    """
    Receives slide content and a theme reference, creates a PowerPoint,
    and returns it as a downloadable file.
    """
    try:
        pptx_file_stream = create_presentation(
            slide_data=presentation_content.slides,
            theme_type=presentation_content.theme_type,
            theme_name=presentation_content.theme_name
        )

        print("Sending PowerPoint presentation to client.")
        return StreamingResponse(
            pptx_file_stream,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename=presentation.pptx"}
        )
    except Exception as e:
        if isinstance(e, HTTPException):
            raise e
        print(f"An unexpected server error occurred: {e}")
        raise HTTPException(status_code=500, detail=f"An unexpected server error occurred: {str(e)}")


@app.get("/")
async def root():
    return {"message": "Welcome to the PDF to Presentation API!"}
