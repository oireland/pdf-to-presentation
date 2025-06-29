# main.py
# Import necessary libraries
import io
import json
import os
import uuid
from typing import List, Literal, Optional

import fitz  # PyMuPDF
import google.generativeai as genai
from dotenv import load_dotenv
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field

# Import theme_config using a relative import
from .theme_config import THEME_CONFIG


# --- Models ---
class Slide(BaseModel):
    title: str = Field(..., description="The title of the slide (maximum 10 words)")
    bullets: Optional[List[str]] = Field(None, description="List of bullet points for the slide")
    text_block: Optional[str] = Field(None, description="A block of text for the slide instead of bullet points")
    image_filename: Optional[str] = Field(None, description="The filename of the image to be included in the slide")


class SlideContent(BaseModel):
    slides: List[Slide] = Field(..., description="List of slides for the presentation")


class PresentationContent(BaseModel):
    slides: List[Slide] = Field(..., description="List of slides for the presentation")
    theme_type: Literal["color", "background"] = Field(..., description="The type of theme to apply.")
    theme_name: str = Field(..., description="The name of the theme to apply (e.g., 'classic_dark', 'blue_gradient').")


# --- Configuration ---
load_dotenv()
app = FastAPI(title="PDF to Presentation API")

if not os.path.exists("extracted_images"):
    os.makedirs("extracted_images")

app.mount("/images", StaticFiles(directory="extracted_images"), name="images")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

try:
    api_key = os.getenv("GOOGLE_API_KEY")
    if not api_key:
        raise ValueError("GOOGLE_API_KEY environment variable not set.")
    genai.configure(api_key=api_key)
except ValueError as e:
    print(f"ERROR: {e}")


# --- Helper Functions ---

def extract_text_and_images_from_pdf(pdf_content: bytes) -> (str, List[dict]):
    full_text = ""
    images = []
    with fitz.open(stream=pdf_content, filetype="pdf") as doc:
        for page_num, page in enumerate(doc):
            full_text += page.get_text()
            for img_index, img_info in enumerate(page.get_images(full=True)):

                # --- START: Definitive Image Processing Logic ---
                xref = img_info[0]
                smask_xref = img_info[1]  # XREF for the soft mask image

                # If there's a soft mask, we need to recombine it with the base image
                if smask_xref > 0:
                    try:
                        print(f"Recombining transparent image on page {page_num + 1} with its mask.")
                        # Get the base image pixmap
                        base_pix = fitz.Pixmap(doc, xref)
                        # Get the mask pixmap
                        mask_pix = fitz.Pixmap(doc, smask_xref)

                        # Create a new pixmap by combining the base and the mask
                        final_pix = fitz.Pixmap(base_pix, mask_pix)

                        # If the base was not RGBA, convert it before cleaning up
                        if base_pix.alpha == 0:
                            base_pix = fitz.Pixmap(fitz.csRGB, base_pix)

                        # Clean up intermediate pixmaps
                        base_pix = None
                        mask_pix = None

                    except Exception as e:
                        print(f"Error combining pixmaps: {e}, falling back to base image.")
                        final_pix = fitz.Pixmap(doc, xref)  # Fallback
                else:
                    # No mask, just get the image directly
                    final_pix = fitz.Pixmap(doc, xref)

                # Now, convert the final pixmap (which has correct transparency) to bytes
                image_bytes = final_pix.tobytes("png")
                final_pix = None  # Clean up
                image_ext = "png"
                # --- END: Definitive Image Processing Logic ---

                image_filename = f"image_{page_num + 1}_{img_index + 1}_{uuid.uuid4()}.{image_ext}"
                image_path = os.path.join("extracted_images", image_filename)

                with open(image_path, "wb") as img_file:
                    img_file.write(image_bytes)

                words = page.get_text("words")
                img_bbox = page.get_image_bbox(img_info)
                context_text = " ".join([w[4] for w in words if fitz.Rect(w[:4]).intersects(img_bbox)])

                images.append({"filename": image_filename, "context": context_text})

    print(f"Successfully extracted {len(full_text)} characters and {len(images)} images.")
    return full_text, images


def generate_slides_content_with_gemini(text: str, images: List[dict], detail_level: int = 2) -> List[Slide]:
    print(f"Generating slide content with Gemini at detail level {detail_level}...")
    model = genai.GenerativeModel('gemini-1.5-flash')
    detail_descriptions = {
        0: "very concise", 1: "concise", 2: "normal",
        3: "detailed", 4: "very detailed"
    }
    bullet_counts = {
        0: "1-2", 1: "2-3", 2: "3-4", 3: "4-5", 4: "5-6"
    }
    detail_description = detail_descriptions.get(detail_level, "normal")
    bullet_count = bullet_counts.get(detail_level, "3-4")
    image_prompts = ""
    if images:
        image_prompts = "\n\nThe following images were extracted from the PDF. Please incorporate them into the slides where they are most relevant, using their context to guide placement. Use each image no more than once:\n"
        for image in images:
            image_prompts += f"- Image: {image['filename']}, Context: {image['context']}\n"
    prompt = f"""
    Based on the following text and images from a report, please generate a {detail_description} summary presentation.
    The output should be a valid JSON object.
    The JSON object must be a single list `[]` containing multiple slide objects {{}}.
    Each slide object must have the following structure:
    {{
        "title": "Slide Title",
        "bullets": ["Point 1", "Point 2", ...],
        "text_block": "A paragraph of text that summarizes the key points...",
        "image_filename": "image_filename.png"
    }}
    About 20% of slides should use text_block instead of bullets for variety.
    For bullet point slides, include approximately {bullet_count} bullet points per slide.
    The level of detail should be {detail_description}.
    {image_prompts}
    Summarize the key points and structure them logically for a presentation.
    Here is the text:
    ---
    {text}
    ---
    """
    try:
        response = model.generate_content(prompt)
        response_text = response.text.strip().replace("```json", "").replace("```", "")
        raw_slide_data = json.loads(response_text)
        if not isinstance(raw_slide_data, list):
            raise ValueError("AI response is not a list.")
        slide_data = [Slide(**slide) for slide in raw_slide_data]
        print("Successfully generated and parsed slide data from Gemini.")
        return slide_data
    except Exception as e:
        print(f"An error occurred during Gemini content generation: {e}")
        raise HTTPException(status_code=500, detail="Failed to generate or parse content from AI.")


def set_slide_background_image(prs, slide, image_path):
    left = top = 0
    width = prs.slide_width
    height = prs.slide_height
    pic = slide.shapes.add_picture(image_path, left, top, width, height)
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def calculate_font_size(text_length: int, min_font=14, max_font=28, min_chars=100, max_chars=500) -> int:
    """Calculates a dynamic font size based on text length."""
    if text_length < min_chars:
        return max_font
    if text_length > max_chars:
        return min_font

    # Linear interpolation between max and min font size
    slope = (min_font - max_font) / (max_chars - min_chars)
    font_size = max_font + slope * (text_length - min_chars)

    return int(font_size)

def create_presentation(slide_data: List[Slide], theme_type: str, theme_name: str) -> io.BytesIO:
    """
    Creates a PowerPoint presentation from structured data using a predefined theme.
    This version uses specific layouts for bullets vs. text_blocks and dynamically
    resizes content placeholders to prevent overlap with images.
    """
    print(f"Creating presentation with new plan: '{theme_name}' (type: {theme_type})")

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # --- Define Slide Layouts ---
    title_and_content_layout = prs.slide_layouts[1]  # For bullet points
    title_only_layout = prs.slide_layouts[5]  # For text blocks

    # --- Get Theme Configuration ---
    theme = None
    if theme_type == "color":
        theme = next((t for t in THEME_CONFIG["colors"] if t["name"] == theme_name), None)
    elif theme_type == "background":
        theme = next((t for t in THEME_CONFIG["backgrounds"] if t["name"] == theme_name), None)

    if not theme:
        print(f"Warning: Theme '{theme_name}' not found.")

    # (Master slide theme application remains the same)
    master = prs.slide_master
    if theme:
        if theme_type == "color":
            fill = master.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(theme["primaryColor"])
            for shape in master.placeholders:
                if shape.has_text_frame:
                    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme["textColor"])
        elif theme_type == "background":
            for shape in master.placeholders:
                if shape.has_text_frame:
                    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme["textColor"])

    # --- Create Slides ---
    for slide_info in slide_data:
        has_image = bool(slide_info.image_filename)

        if slide_info.bullets:
            slide = prs.slides.add_slide(title_and_content_layout)
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]

            # Resize body placeholder if there is an image
            if has_image:
                body_shape.width = Inches(7.5)
                body_shape.top = Inches(1.5)
                body_shape.left = Inches(0.5)


            # Populate bullets
            tf = body_shape.text_frame
            tf.clear()

            total_chars = sum(len(s) for s in slide_info.bullets)
            font_size = calculate_font_size(total_chars)

            for bullet_text in slide_info.bullets:
                p = tf.add_paragraph()
                p.text = bullet_text
                p.font.size = Pt(font_size)
                p.level = 0
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        elif slide_info.text_block:
            slide = prs.slides.add_slide(title_only_layout)
            title_shape = slide.shapes.title

            # Define dimensions for the text box
            left = Inches(0.5)
            top = Inches(1.5)
            height = Inches(7.0)
            # Adjust width if there is an image
            width = Inches(7.5) if has_image else Inches(15)

            # Add and populate the text box
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = slide_info.text_block

            font_size = calculate_font_size(len(slide_info.text_block))
            p.font.size = Pt(font_size)

            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            tf.word_wrap = True

        else:  # Title only slide
            slide = prs.slides.add_slide(title_only_layout)
            title_shape = slide.shapes.title

        # --- Set Title ---
        if title_shape:
            title_shape.text = slide_info.title
            if theme:
                if theme_type == "color":
                    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme["secondaryColor"])
                elif theme_type == "background":
                    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme["titleColor"])

        # --- Add Image (if it exists) ---
        if has_image:
            image_path = os.path.join("extracted_images", slide_info.image_filename)
            if os.path.exists(image_path):
                # Place image on the right side of the slide
                slide.shapes.add_picture(image_path, Inches(8.25), Inches(1.5), width=Inches(7.25))

        # --- Set Slide Background (for image themes) ---
        if theme and theme_type == "background":
            image_path = os.path.join(os.path.dirname(__file__), "backgrounds", theme["image_filename"])
            if os.path.exists(image_path):
                slide.background.fill.solid()
                slide.background.fill.picture(image_path)
            else:
                print(f"Warning: Background image not found at {image_path}")

    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    return pptx_stream


# --- API Endpoints ---
@app.post("/api/generate-slide-content", response_model=SlideContent)
async def generate_slide_content(file: UploadFile = File(...), detail_level: int = 2):
    if file.content_type != "application/pdf":
        raise HTTPException(status_code=400, detail="Invalid file type. Please upload a PDF.")
    try:
        pdf_contents = await file.read()
        text_content, images = extract_text_and_images_from_pdf(pdf_contents)
        slides = generate_slides_content_with_gemini(text_content, images, detail_level)
        return SlideContent(slides=slides)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/generate-presentation")
async def generate_presentation_endpoint(presentation_content: PresentationContent):
    try:
        pptx_file_stream = create_presentation(
            slide_data=presentation_content.slides,
            theme_type=presentation_content.theme_type,
            theme_name=presentation_content.theme_name
        )
        return StreamingResponse(
            pptx_file_stream,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=presentation.pptx"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/images/{image_filename}")
async def get_image(image_filename: str):
    image_path = os.path.join("extracted_images", image_filename)
    if os.path.exists(image_path):
        return FileResponse(image_path)
    raise HTTPException(status_code=404, detail="Image not found")


@app.get("/")
async def root():
    return {"message": "Welcome to the PDF to Presentation API!"}