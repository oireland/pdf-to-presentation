# main.py
# Import necessary libraries
import os
import fitz  # PyMuPDF
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import JSONResponse
from fastapi.middleware.cors import CORSMiddleware
import google.generativeai as genai
from pptx import Presentation
import json
import io

# --- Configuration ---
# Create a FastAPI app instance
app = FastAPI(title="PDF to Presentation API")

# Configure CORS (Cross-Origin Resource Sharing)
# This allows your React frontend to communicate with this backend.
# IMPORTANT: For production, you should restrict the origins to your frontend's domain.
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allows all origins for now
    allow_credentials=True,
    allow_methods=["*"],  # Allows all methods
    allow_headers=["*"],  # Allows all headers
)

# Configure the Google Gemini API key
try:
    genai.configure(api_key=os.environ["GOOGLE_API_KEY"])
except KeyError:
    # This provides a helpful error message if the key is not set.
    print("ERROR: GOOGLE_API_KEY environment variable not set.")
    # In a real app, you might want to exit or disable AI features.


# --- Helper Functions (We will build these out later) ---

def extract_text_from_pdf(file_stream):
    """Extracts text content from a PDF file stream."""
    # This is a placeholder for now. We will implement this in the next step.
    print("Placeholder: Extracting text from PDF...")
    return "This is placeholder text from the PDF."


def generate_slides_content_with_gemini(text):
    """Uses Gemini to generate presentation content from text."""
    # This is a placeholder for now. We will implement this in the next step.
    print("Placeholder: Generating content with Gemini...")
    # This is the kind of structured data we expect from the AI
    slide_data = [
        {"title": "Slide 1 Title", "bullets": ["Point 1", "Point 2"]},
        {"title": "Slide 2 Title", "bullets": ["Point A", "Point B"]},
    ]
    return slide_data


def create_presentation(slide_data):
    """Creates a PowerPoint presentation from structured data."""
    # This is a placeholder for now. We will implement this in the next step.
    print("Placeholder: Creating PowerPoint file...")
    # Create an in-memory presentation
    prs = Presentation()
    for slide_info in slide_data:
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        title.text = slide_info.get("title", "No Title")

        # Add bullet points
        tf = body.text_frame
        tf.clear()  # Clear existing text (like default prompts)
        for bullet in slide_info.get("bullets", []):
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

    # Save the presentation to a byte stream
    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)  # Rewind the stream to the beginning
    return pptx_stream


# --- API Endpoint ---

@app.post("/api/generate-presentation")
async def generate_presentation_endpoint(file: UploadFile = File(...)):
    """
    This endpoint receives a PDF file, processes it, and returns a
    PowerPoint presentation.
    """
    print(f"Received file: {file.filename} ({file.content_type})")

    # Basic validation for PDF file type
    if file.content_type != "application/pdf":
        raise HTTPException(status_code=400, detail="Invalid file type. Please upload a PDF.")

    try:
        # --- Core Logic (will be fully implemented later) ---
        # 1. Extract text from PDF
        # Note: We will replace this with the real implementation
        text_content = extract_text_from_pdf(None)  # Passing None for now

        # 2. Generate slide content with Gemini
        # Note: We will replace this with the real implementation
        slides = generate_slides_content_with_gemini(text_content)

        # 3. Create the .pptx file in memory
        pptx_file_stream = create_presentation(slides)

        # In a real implementation, we would send this back.
        # For now, we'll return a success message.
        print("Successfully created presentation stream.")

        return JSONResponse(
            status_code=200,
            content={"message": "Presentation generated successfully! (This is a placeholder response)"}
        )

    except Exception as e:
        # Generic error handling
        print(f"An error occurred: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/")
def read_root():
    """A simple root endpoint to confirm the server is running."""
    return {"message": "Welcome to the Presentation Generator API!"}

